import json
import os
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pathlib import Path
from PIL import Image, ImageTk
from google import genai
from google.genai import types
from pydantic import BaseModel
from config import config_ini
import logging
from get_prompt import get_system_instructions
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from pptx.util import Inches, Pt
from math import ceil, floor
from datetime import datetime

# ロギング設定
output_file = config_ini.get("LOGGING", "log_file", fallback="app.log")
encoding = config_ini.get("LOGGING", "encoding", fallback="utf-8")
level_name = config_ini.get("LOGGING", "log-level", fallback="INFO").upper()
level = getattr(logging, level_name, logging.INFO)
log_format = config_ini.get(
    "LOGGING",
    "format",
    fallback="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
)

# ルートロガーの設定
logging.basicConfig(
    level=level,  # ここが重要：ルートロガーのレベルを設定
    format=log_format,
    handlers=[
        logging.FileHandler(output_file, encoding=encoding),
        logging.StreamHandler(),
    ],
)

logger = logging.getLogger(__name__)

# このファイル（main.py）がある場所を取得
BASE_DIR = Path(__file__).resolve().parent


class ImageTextboxApp:
    def __init__(self, root, config_ini):
        self.root = root
        self.root.title("画像プレビューアプリケーション")
        self.config_ini = config_ini
        self.output_dir = self.config_ini.get(
            "PPTX_SETTINGS", "output_dir", fallback="pptx_output"
        )
        # 絶対パスに変換
        self.output_dir = BASE_DIR / self.output_dir

        self.root.geometry(
            config_ini.get("GUI_SETTINGS", "window_size", fallback="1170x450")
        )

        self.apiKey = config_ini.get("GEMINI", "api_key", fallback="")
        if not self.apiKey:
            logger.warning(
                "GEMINI APIキーが設定されていません。APIキーを設定してください。"
            )
            raise ValueError("GEMINI APIキーが設定されていません。")
        # Gemini APIクライアントの初期化
        self.generate_client = genai.Client(api_key=self.apiKey)

        # アップロードされた画像のパスを保存
        self.uploaded_images = []

        self.gemini_model = config_ini.get(
            "GEMINI", "model", fallback="gemini-2.5-flash"
        )
        try:
            self.system_instruction = (
                get_system_instructions()
                or "You are a helpful assistant that extracts text from images."
            )
        except FileNotFoundError as fnf_error:
            logger.exception("System instruction file error")
            self.system_instruction = (
                "You are a helpful assistant that extracts text from images."
            )

        # メインコンテナ
        self.setup_ui()

    def setup_ui(self):
        # メインコンテンツエリア（パーン分割）
        self.paned_window = ttk.PanedWindow(self.root, orient=tk.HORIZONTAL)
        self.paned_window.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # 左側パネル: コンポーネント（複数階層）
        self.setup_left_panel()

        # 右側パネル: テキストボックス（画像プレビュー）
        self.setup_right_panel()

        logger.info("UI setup complete")

    def setup_left_panel(self):
        # 左側フレーム
        left_frame = ttk.Frame(self.paned_window, width=220)
        self.paned_window.add(left_frame, weight=1)

        # ラベル
        label = ttk.Label(
            left_frame, text="Image to Textbox", font=("Arial", 10, "bold")
        )
        label.pack(pady=5, anchor=tk.W, padx=5)

        # pptxファイル名入力フレーム
        file_frame = ttk.Frame(left_frame)
        file_frame.pack(fill=tk.X, padx=5, pady=5)

        self.file_name = tk.StringVar()
        ttk.Label(file_frame, text="ファイル名:").pack(side=tk.LEFT, padx=(0, 5))
        self.file_name_entry = ttk.Entry(
            file_frame, textvariable=self.file_name, width=25
        )
        self.file_name_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)

        # モデル名表示フレーム
        model_frame = ttk.Frame(left_frame)
        model_frame.pack(fill=tk.X, padx=5, pady=5)

        ttk.Label(model_frame, text="モデル名:").pack(side=tk.LEFT, padx=(0, 5))
        self.model_name_label = ttk.Label(
            model_frame, text=self.gemini_model, relief=tk.SUNKEN, width=25, anchor=tk.W
        )
        self.model_name_label.pack(side=tk.LEFT, fill=tk.X, expand=True)

        # ファイルアップロードボタン
        upload_frame = ttk.Frame(left_frame)
        upload_frame.pack(fill=tk.BOTH, padx=5, pady=5)

        ttk.Button(
            upload_frame, text="ファイルをアップロード", command=self.on_file_upload
        ).pack(side=tk.LEFT, fill=tk.X, expand=True)
        ttk.Button(upload_frame, text="リセット", command=self.on_reset).pack(
            side=tk.LEFT, fill=tk.X, expand=True
        )

        # ファイル名一覧フレーム
        file_list_frame = ttk.LabelFrame(left_frame, text="ファイル名一覧", padding=5)
        file_list_frame.pack(fill=tk.BOTH, expand=False, padx=5, pady=5)

        # ファイル名一覧用のリストボックス
        list_scrollbar = ttk.Scrollbar(file_list_frame, orient=tk.VERTICAL)
        list_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        self.file_listbox = tk.Listbox(
            file_list_frame, yscrollcommand=list_scrollbar.set, height=8
        )
        self.file_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        list_scrollbar.config(command=self.file_listbox.yview)

        # コントロールボタンフレーム（開始/停止）
        control_frame = ttk.Frame(left_frame)
        control_frame.pack(fill=tk.Y, padx=5, pady=5)

        self.start_button = ttk.Button(
            control_frame, text="開始", width=12, command=self.on_start
        )
        self.start_button.pack(side=tk.LEFT, padx=2)

        self.stop_button = ttk.Button(
            control_frame,
            text="停止",
            width=12,
            command=self.on_stop,
            state=tk.DISABLED,
        )
        self.stop_button.pack(side=tk.LEFT, padx=2)

        # ステータス表示フレーム
        status_frame = ttk.LabelFrame(left_frame, text="ステータス", padding=5)
        status_frame.pack(
            fill=tk.BOTH,
            padx=5,
            pady=5,
            expand=True,
        )
        self.status_display = ttk.Label(
            status_frame, text="準備完了", anchor=tk.W, font=("Arial", 12)
        )
        self.status_display.pack(fill=tk.X)

    def setup_right_panel(self):
        # 右側フレーム
        right_frame = ttk.Frame(self.paned_window)
        self.paned_window.add(right_frame, weight=3)

        # ラベル
        label = ttk.Label(
            right_frame,
            text="画像プレビュー（すべての画像）",
            font=("Arial", 10, "bold"),
        )
        label.pack(pady=5, anchor=tk.W, padx=4)

        # 画像プレビューエリア（スクロール可能なキャンバス）
        canvas_frame = ttk.Frame(right_frame)
        canvas_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # スクロールバー
        v_scrollbar = ttk.Scrollbar(canvas_frame, orient=tk.VERTICAL)
        v_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        h_scrollbar = ttk.Scrollbar(canvas_frame, orient=tk.HORIZONTAL)
        h_scrollbar.pack(side=tk.BOTTOM, fill=tk.X)

        # キャンバス
        self.image_canvas = tk.Canvas(
            canvas_frame,
            bg="white",
            yscrollcommand=v_scrollbar.set,
            xscrollcommand=h_scrollbar.set,
        )
        self.image_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        v_scrollbar.config(command=self.image_canvas.yview)
        h_scrollbar.config(command=self.image_canvas.xview)

        # キャンバス内にフレームを作成
        self.images_frame = ttk.Frame(self.image_canvas)
        self.canvas_window = self.image_canvas.create_window(
            (0, 0), window=self.images_frame, anchor=tk.NW
        )

        # フレームのサイズが変更されたときにスクロール領域を更新
        self.images_frame.bind("<Configure>", self.on_frame_configure)

        # プレースホルダーラベル
        self.placeholder_label = ttk.Label(
            self.images_frame,
            text="画像ファイルをアップロードしてください",
            font=("Arial", 12),
            anchor=tk.CENTER,
        )
        self.placeholder_label.pack(pady=50)

    def on_frame_configure(self, event=None):
        """キャンバスのスクロール領域を更新"""
        self.image_canvas.configure(scrollregion=self.image_canvas.bbox("all"))

    def on_create_folder(self):
        folder_name = self.file_name_entry.get().strip()
        if folder_name:
            self.status_display.config(text=f"フォルダ '{folder_name}' を作成中...")
            # ここで実際のフォルダ作成処理を行う
            try:
                # 例: Path(folder_name).mkdir(exist_ok=True)
                messagebox.showinfo("成功", f"フォルダ '{folder_name}' を作成しました")
                self.status_display.config(text="準備完了")
            except Exception as e:
                messagebox.showerror("エラー", f"フォルダ作成に失敗しました: {e}")
                logger.exception("フォルダ作成エラー")
                self.status_display.config(text="エラー")
        else:
            messagebox.showwarning("警告", "フォルダ名を入力してください")

    def on_file_upload(self):
        file_paths = filedialog.askopenfilenames(
            title="ファイルを選択",
            filetypes=[
                ("画像ファイル (JPEG/PNG)", "*.jpg *.jpeg *.png"),
            ],
        )
        if file_paths:
            for file_path in file_paths:
                file_name = Path(file_path).name
                # リストボックスに追加（重複チェック）
                if file_name not in self.file_listbox.get(0, tk.END):
                    self.file_listbox.insert(tk.END, file_name)
                    # 画像パスを保存
                    self.uploaded_images.append(file_path)

            self.status_display.config(
                text=f"{len(file_paths)}個のファイルをアップロードしました"
            )

            # 画像を表示
            self.display_images()

    def on_reset(self):
        """リセットボタンの処理"""
        # ファイルリストをクリア
        self.file_listbox.delete(0, tk.END)
        self.uploaded_images.clear()

        # 画像表示エリアをクリア
        for widget in self.images_frame.winfo_children():
            widget.destroy()

        self.on_frame_configure()

        # プレースホルダーラベルを再表示
        self.placeholder_label = ttk.Label(
            self.images_frame,
            text="画像ファイルをアップロードしてください",
            font=("Arial", 12),
        )
        self.placeholder_label.pack(pady=50)

        self.status_display.config(text="リセット完了")

    def display_images(self):
        """アップロードされた画像をすべて表示（2列レイアウト）"""
        # プレースホルダーを削除
        if hasattr(self, "placeholder_label"):
            self.placeholder_label.destroy()

        # 既存の画像ウィジェットをクリア
        for widget in self.images_frame.winfo_children():
            widget.destroy()

        # 画像参照を保持するリスト（ガベージコレクション防止）
        self.image_references = []

        # 2列レイアウト用の行フレーム
        current_row_frame = None

        # 各画像を表示
        for idx, img_path in enumerate(self.uploaded_images):
            try:
                # 2列ごとに新しい行フレームを作成
                if idx % 2 == 0:
                    current_row_frame = ttk.Frame(self.images_frame)
                    current_row_frame.pack(fill=tk.X, pady=5)

                # 画像を読み込み
                with Image.open(img_path) as img:

                    # サムネイルサイズに縮小（アスペクト比を維持）
                    img.thumbnail((325, 325), Image.Resampling.LANCZOS)

                    # PhotoImageに変換
                    photo = ImageTk.PhotoImage(img)
                self.image_references.append(photo)

                # フレームを作成（2列配置）
                img_container = ttk.Frame(
                    current_row_frame, relief=tk.RIDGE, borderwidth=2
                )
                img_container.pack(side=tk.LEFT, pady=5, padx=5, expand=True)

                # ファイル名ラベル
                name_label = ttk.Label(
                    img_container,
                    text=Path(img_path).name,
                    font=("Arial", 9, "bold"),
                    wraplength=330,
                )
                name_label.pack(pady=5, padx=5)

                # 画像ラベル
                img_label = tk.Label(img_container, image=photo, bg="white")
                img_label.pack(pady=5, padx=5)

            except Exception as e:
                # エラー時は警告を表示
                if idx % 2 == 0:
                    error_row = ttk.Frame(self.images_frame)
                    error_row.pack(fill=tk.X, pady=5)
                    current_row_frame = error_row

                error_label = ttk.Label(
                    current_row_frame,
                    text=f"エラー: {Path(img_path).name} - {str(e)}",
                    foreground="red",
                )
                error_label.pack(side=tk.LEFT, pady=5, padx=5)

    # gemini apiのファイルAPIを使った画像のアップロード
    def file_upload_to_gemini(self):
        """例外を親関数に伝播させる"""
        if not self.uploaded_images:
            logger.warning("アップロードする画像がありません")
            raise ValueError("アップロードする画像がありません")

        # 並列アップロード（最大10スレッド）
        task_list = []
        total_files = len(self.uploaded_images)
        max_workers = min(10, total_files or 1)

        def upload_file(file_path):
            client = genai.Client(api_key=self.apiKey)
            return client.files.upload(file=file_path)

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            for idx, result in enumerate(
                executor.map(upload_file, self.uploaded_images), 1
            ):
                task_list.append(result)
                logger.info(f"Uploaded {idx}/{total_files} files to Gemini")
                self.status_display.config(
                    text=f"アップロード中... {idx}/{total_files} files"
                )

        logger.info(f"Total uploaded: {len(task_list)} files")
        return task_list

    def _delete_file(self, file_id):
        client = genai.Client(api_key=self.apiKey)
        client.files.delete(name=file_id.name)

    def extract_text(self, files):
        """例外を親関数に伝播させる"""

        class figure_token(BaseModel):
            figure_name: str
            token: list[str]

        if not files:
            logger.warning("テキスト抽出のためのファイルがありません")
            raise ValueError("テキスト抽出のためのファイルがありません")
        logger.info("Starting text extraction")
        self.status_display.config(text="テキスト抽出中...")

        response = self.generate_client.models.generate_content(
            model=self.gemini_model,
            config=types.GenerateContentConfig(
                system_instruction=self.system_instruction,
                response_mime_type="application/json",
                response_schema=list[figure_token],
            ),
            contents=[*files, "添付した画像について処理を行ってください。"],
        )
        max_workers = min(10, len(files) or 1)

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            for idx, _ in enumerate(executor.map(self._delete_file, files)):
                logger.info(f"Deleted {idx}/{len(files)} files from Gemini")

        # None または text欠如を検出
        if not response or getattr(response, "text", None) is None:
            raise ValueError("No response text received from Gemini API")

        # 空文字列を検出
        if not response.text:
            raise ValueError("Empty response text received from Gemini API")

        json_response = json.loads(response.text)  # 例外はここで発生（親に伝播）
        logger.info("Text extraction successful")
        return json_response

    def generate_pptx(self, gemini_response):
        prs = Presentation()
        # 設定値をロード
        font_name = self.config_ini.get("PPTX_SETTINGS", "font_name", fallback="Arial")
        font_size = self.config_ini.getint("PPTX_SETTINGS", "font_size", fallback=14)
        layout_num = self.config_ini.getint("PPTX_SETTINGS", "layout_num", fallback=6)
        char_width_in = self.config_ini.getfloat(
            "PPTX_SETTINGS", "char_width_in", fallback=0.097
        )
        min_w_in = self.config_ini.getfloat("PPTX_SETTINGS", "min_w_in", fallback=0.45)
        min_h_in = self.config_ini.getfloat("PPTX_SETTINGS", "min_h_in", fallback=0.30)
        wrap_padding_in = self.config_ini.getfloat(
            "PPTX_SETTINGS", "wrap_padding_in", fallback=0.20
        )

        # Heading box (single full-width box at the top)
        margin_l = self.config_ini.getfloat("PPTX_SETTINGS", "margin_l", fallback=0.4)
        margin_r = self.config_ini.getfloat("PPTX_SETTINGS", "margin_r", fallback=0.4)
        margin_t = self.config_ini.getfloat("PPTX_SETTINGS", "margin_t", fallback=0.5)
        margin_b = self.config_ini.getfloat("PPTX_SETTINGS", "margin_b", fallback=0.4)
        heading_h = self.config_ini.getfloat("PPTX_SETTINGS", "heading_h", fallback=0.4)
        line_height_in = 1.3 * (font_size / 72.0)

        def add_token_grid_slide(prs, title, token_list, cols=4):
            layouts = prs.slide_layouts
            idx = layout_num if 0 <= layout_num < len(layouts) else 6

            slide = prs.slides.add_slide(layouts[idx])  # blank layout

            # Page geometry
            page_w = prs.slide_width / 914400.0  # EMU -> inches
            page_h = prs.slide_height / 914400.0

            # Title
            title_box = slide.shapes.add_textbox(
                Inches(margin_l),
                Inches(margin_t - 0.1),
                Inches(page_w - margin_l - margin_r),
                Inches(heading_h),
            )
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"Tokens from panel {title}"
            run.font.size = Pt(font_size)

            # Grid region
            grid_top = margin_t + heading_h + 0.1
            grid_left = margin_l
            grid_w = page_w - margin_l - margin_r
            grid_h = page_h - grid_top - margin_b

            n = len(token_list)
            rows = ceil(n / cols) if n else 1
            cell_w_in = grid_w / cols
            cell_h_in = grid_h / rows

            for idx, token in enumerate(token_list):
                r = idx // cols
                c = idx % cols

                # Size optimization rules
                w_in = max(
                    min_w_in,
                    min(cell_w_in, char_width_in * len(token) + wrap_padding_in),
                )
                max_chars = max(1, floor((w_in - wrap_padding_in) / char_width_in))
                lines = max(1, ceil(len(token) / max_chars))
                h_in = max(min_h_in, min(0.9 * cell_h_in, lines * line_height_in))

                # Position within the cell (top-left, with small padding)
                cell_x = grid_left + c * cell_w_in
                cell_y = grid_top + r * cell_h_in
                pad = 0.05
                left = cell_x + pad
                top = cell_y + pad

                box = slide.shapes.add_textbox(
                    Inches(left), Inches(top), Inches(w_in), Inches(h_in)
                )
                tf = box.text_frame
                tf.clear()  # required by spec
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = token
                run.font.name = font_name
                run.font.size = Pt(font_size)

            return slide

        for figure in gemini_response:
            add_token_grid_slide(
                prs,
                figure.get("figure_name", "Unknown"),
                figure.get("token", []),
                cols=4,
            )

        # 保存
        safe_name = self.file_name.get().strip()
        if not safe_name:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            pptx_filename = f"output_{timestamp}.pptx"
        else:
            # パストラバーサル対策: ベース名のみを使用
            safe_stem = Path(safe_name).stem
            safe_basename = Path(safe_stem).name  # ディレクトリ分を除去
            pptx_filename = f"{safe_basename}.pptx"

        output_path = self.output_dir / pptx_filename

        # 出力ディレクトリを確実に作成
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # パストラバーサル検証
        try:
            output_path.resolve().relative_to(self.output_dir.resolve())
        except ValueError:
            logger.exception("パストラバーサルの試行を検出しました")
            raise ValueError("無効なファイル名が指定されました")

        try:
            prs.save(output_path)
            logger.info("PPTXファイルを保存しました: %s", output_path)
        except Exception:
            logger.exception("PPTXファイルの保存中にエラーが発生しました")
            raise

        return output_path

    def on_start(self):
        """開始ボタンの処理"""
        if self.file_listbox.size() == 0:
            messagebox.showwarning("警告", "ファイルをアップロードしてください")
            logger.warning(
                "ファイルがアップロードされていません。処理を開始できません。"
            )
            return

        self.status_display.config(text="処理を開始しました")
        self.start_button.config(state=tk.DISABLED)
        self.stop_button.config(state=tk.NORMAL)

        # ここで実際の処理を開始
        try:
            logger.info("処理を開始しました。")
            # 出力ディレクトリの存在確認
            self.output_dir.mkdir(parents=True, exist_ok=True)
        except ValueError as ve:
            messagebox.showerror("エラー", f"処理中にエラーが発生しました: {ve}")
            logger.exception("ValueError during processing")
        except Exception as e:
            messagebox.showerror(
                "エラー", f"処理中に予期しないエラーが発生しました: {e}"
            )
            logger.exception("Unexpected error during processing")
        finally:
            self.on_finish()

    def on_stop(self):
        """停止ボタンの処理"""
        self.status_display.config(text="処理を停止しました")
        self.start_button.config(state=tk.NORMAL)
        self.stop_button.config(state=tk.DISABLED)

        # ここで実際の処理を停止
        messagebox.showinfo("停止", "処理を停止しました")
        self.status_display.config(text="準備完了")

    def on_finish(self):
        """処理完了時の共通処理"""
        self.start_button.config(state=tk.NORMAL)
        self.stop_button.config(state=tk.DISABLED)
        self.status_display.config(text="準備完了")
        messagebox.showinfo("完了", "処理が完了しました")


def main():
    root = tk.Tk()
    icon_name = config_ini.get("GUI_SETTINGS", "icon_name", fallback="favicon.ico")
    icon_path = os.path.join(os.path.dirname(__file__), "config", icon_name)
    try:
        root.iconbitmap(default=icon_path)
    except Exception as e:
        logger.exception("アイコンの設定に失敗しました")

    try:
        ImageTextboxApp(root, config_ini)
    except ValueError as ve:
        logger.exception("アプリケーションの初期化に失敗しました")
        messagebox.showerror("エラー", f"アプリケーションの初期化に失敗しました: {ve}")
        root.destroy()
        return

    root.mainloop()


if __name__ == "__main__":
    main()
